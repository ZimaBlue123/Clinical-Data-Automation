#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Intelligent PPT Restructuring & Narrative Flow Automation (ppt_engine)

按「终极版 PPT 叙事剧本」做去重与精选：
1. 仅保留剧本中列出的幻灯片，其余删除（信息冗余或未入选）；
2. 严格按五章顺序 + 每章内剧本顺序输出；
3. 每章前插入过渡页（深绿底、居中大标题）；
4. 若多页匹配同一剧本条目，保留信息量更大（shape_count 更大）的一页。
"""

import os
import re
from typing import List, Dict, Tuple, Optional

# 复用 merge_ppt 的合并与复制能力（input/ 读源文件，output/ 写结果）
from merge_ppt import (
    get_merge_dir,
    get_output_dir,
    auto_discover_ppts,
    build_slide_infos,
    clone_slide_into_presentation,
    standardize_fonts,
    SlideInfo,
)
from merge_ppt import extract_slide_text
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# =============================================================================
# 终极版叙事剧本：五章 + 每章保留的幻灯片（按标题片段匹配，顺序即播放顺序）
# =============================================================================

# (章节ID, 章节显示名, 该章内「槽位」列表；每槽位为短语列表，任一匹配即进入该槽)
CURATED_SCRIPT: List[Tuple[str, str, List[List[str]]]] = [
    ("1", "研究背景与设计 (Context)", [
        ["研究设计与安全集分布", "Study Design", "安全集", "Enrollment", "design", "flowchart"],
    ]),
    ("2", "核心安全性结论 (Top-line Conclusion)", [
        ["总体安全性概览", "Overall Safety", "Top-line", "400", "非临床剂量", "Safety Summary"],
    ]),
    ("3", "详细安全性特征 (Detailed Safety Profile)", [
        ["0-30天相关AE", "0-30", "30-day", "TEAE"],
        ["MedDRA PT", "PT分类", "发热", "肌痛", "MedDRA", "preferred term"],
        ["0-14天严重程度", "Grade 3", "严重程度", "severity", "0-14"],
        ["持续时间", "Duration", "≥8天", "1.0%", "8 days", "duration"],
        ["佐剂剂量反应", "佐剂", "dose response", "adjuvant"],
    ]),
    ("4", "竞品横向对标 (Benchmark vs Shingrix)", [
        ["Shingrix 临床", "Shingrix", "GSK", "clinical"],
        ["反应原性痛点", "reactogenicity", "GSK", "pain"],
        ["接种部位疼痛", "疼痛对比", "injection site", "pain at injection"],
        ["总体耐受性优势", "耐受性", "benchmark", "tolerability"],
    ]),
    ("5", "总结 (Conclusion)", [
        ["总结与临床价值", "Conclusion", "依从性", "临床价值", "compliance", "Summary"],
    ]),
]

# 当标题/正文匹配不到时，按「原始合并顺序下的 Slide 序号」(1-based) 直接指定保留页。
CURATED_SLIDE_NUMBERS: List[Tuple[str, List[int]]] = [
    ("1", [24]), ("2", [25]), ("3", [19, 26, 18, 27, 21]), ("4", [13, 14, 29, 30]), ("5", [31]),
]

# =============================================================================
# 严格去重与叙事顺序：「编辑」逻辑 (SLIDE_BLUEPRINT)
# 每 Slot 仅保留一张优胜 Slide；多张匹配时按 tie_breaker 决出唯一胜者。
# =============================================================================

# match 键: title_all / title_any / text_all / content_any（在标题+正文中）
# tie: first / max_shape_count / body_contains_n64_n32 / has_table / max_text_len
SLIDE_BLUEPRINT: List[Dict] = [
    {"slot": 1, "chapter_id": "1", "chapter_title": "研究背景 (Context)",
     "match": {"title_all": ["重组带状疱疹疫苗", "阶段性分析"]}, "tie": "first"},
    {"slot": 2, "chapter_id": "1", "chapter_title": "研究背景 (Context)",
     "match": {"title_any": ["设计", "Safety Set", "安全集"]}, "tie": "body_contains_n64_n32"},
    {"slot": 3, "chapter_id": "2", "chapter_title": "核心结论 (Top-line Safety)",
     "match": {"text_all": ["0", "98%", "61%"]}, "tie": "first"},
    {"slot": 4, "chapter_id": "3", "chapter_title": "安全性特征详述 (Detailed Profile)",
     "match": {"title_all": ["0-30天", "相关"]}, "tie": "max_shape_count"},
    {"slot": 5, "chapter_id": "3", "chapter_title": "安全性特征详述 (Detailed Profile)",
     "match": {"title_any": ["MedDRA", "分类"]}, "tie": "has_table"},
    {"slot": 6, "chapter_id": "3", "chapter_title": "安全性特征详述 (Detailed Profile)",
     "match": {"title_any": ["持续时间", "Duration"]}, "tie": "first"},
    {"slot": 7, "chapter_id": "3", "chapter_title": "安全性特征详述 (Detailed Profile)",
     "match": {"title_any": ["佐剂", "Adjuvant"]}, "tie": "first"},
    {"slot": 8, "chapter_id": "4", "chapter_title": "竞品对标 (Benchmark vs Shingrix)",
     "match": {"title_any": ["疼痛"], "title_also_any": ["Shingrix", "对比"]}, "tie": "max_shape_count"},
    {"slot": 9, "chapter_id": "4", "chapter_title": "竞品对标 (Benchmark vs Shingrix)",
     "match": {"title_any": ["优势", "分析"]}, "tie": "has_table"},
    {"slot": 10, "chapter_id": "5", "chapter_title": "总结 (Conclusion)",
     "match": {"title_any": ["总结", "临床意义"], "content_any": ["依从性", "II期"]}, "tie": "first"},
]
# 当 match 无法命中时（如标题/正文为空），按 1-based Slide 序号回退取页。每 Slot 可填多个备选序号。
BLUEPRINT_FALLBACK_SLIDES: Dict[int, List[int]] = {
    1: [1],          # 标题页：通常为第 1 张
    2: [20, 2],      # 详细研究设计
    3: [3],
    4: [4, 15],
    5: [9, 18],
    6: [8],
    7: [10],
    8: [13, 29],
    9: [14, 30],
    10: [12, 31],
}


def _slide_has_table(slide) -> bool:
    """判断幻灯片是否包含表格。"""
    try:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                return True
    except Exception:
        pass
    return False


def _get_slide_title(slide) -> str:
    """从单页 Slide 提取标题：优先占位符标题，否则取正文首行。"""
    try:
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip()
    except Exception:
        pass
    full = extract_slide_text(slide)
    if not full:
        return ""
    first_line = full.split("\n")[0].strip()
    return first_line[:200] if first_line else ""


def _slide_matches_blueprint_slot(s: SlideInfo, spec: Dict) -> bool:
    """判断 slide 是否匹配该 Slot 的 match 规则。"""
    title = (getattr(s, "slide_title", "") or "").strip()
    text = ((getattr(s, "slide_title", "") or "") + "\n" + (s.text_content or "")).strip()
    m = spec.get("match", {})
    if not m:
        return False
    # title_all: 标题中必须包含全部
    for key in ("title_all",):
        if key not in m:
            continue
        for phrase in m[key]:
            if phrase.strip() not in title:
                return False
    # title_any: 标题中包含任一
    if "title_any" in m:
        if not any(p.strip() in title for p in m["title_any"] if p.strip()):
            return False
    # title_also_any: 标题中还需包含另一组任一（与 title_any 同时满足）
    if "title_also_any" in m:
        if not any(p.strip() in title for p in m["title_also_any"] if p.strip()):
            return False
    # text_all: 全文（标题+正文）包含全部
    if "text_all" in m:
        for phrase in m["text_all"]:
            if phrase.strip() not in text:
                return False
    # content_any: 全文包含任一
    if "content_any" in m:
        if not any(p.strip() in text for p in m["content_any"] if p.strip()):
            return False
    return True


def _tie_break(candidates: List[SlideInfo], tie: str, slot_id: int, log: List[str]) -> Optional[SlideInfo]:
    """从候选列表中按 tie 规则选出唯一胜者，并写日志。"""
    if not candidates:
        return None
    if tie == "first":
        log.append(f"Slot {slot_id}: Single or first match kept (global_id={candidates[0].global_id}).")
        return candidates[0]
    if tie == "max_shape_count":
        best = max(candidates, key=lambda x: x.shape_count)
        if len(candidates) > 1:
            others = [c.global_id for c in candidates if c is not best]
            log.append(f"Slot {slot_id}: Conflict. Kept slide global_id={best.global_id} (shape_count={best.shape_count}), discarded {others}.")
        else:
            log.append(f"Slot {slot_id}: Kept global_id={best.global_id}.")
        return best
    if tie == "body_contains_n64_n32":
        with_n = [c for c in candidates if ("N=64" in (c.text_content or "") or "N=32" in (c.text_content or ""))]
        if with_n:
            chosen = with_n[0]
            discarded = [c.global_id for c in candidates if c is not chosen]
            log.append(f"Slot {slot_id} (Design): Conflict. Kept slide with 'N=64'/'N=32' (global_id={chosen.global_id}), discarded {discarded}.")
            return chosen
        best = max(candidates, key=lambda x: x.shape_count)
        log.append(f"Slot {slot_id}: No 'N=64'/'N=32' found. Kept global_id={best.global_id} by shape_count.")
        return best
    if tie == "has_table":
        with_table = [c for c in candidates if getattr(c, "has_table", False)]
        if with_table:
            chosen = with_table[0]
            if len(candidates) > 1:
                log.append(f"Slot {slot_id}: Conflict. Kept slide with table (global_id={chosen.global_id}), discarded others (no table).")
            else:
                log.append(f"Slot {slot_id}: Kept global_id={chosen.global_id} (has table).")
            return chosen
        if candidates:
            log.append(f"Slot {slot_id}: No slide with table. Kept first (global_id={candidates[0].global_id}).")
            return candidates[0]
        return None
    if tie == "max_text_len":
        best = max(candidates, key=lambda x: len(x.text_content or ""))
        if len(candidates) > 1:
            log.append(f"Slot {slot_id}: Conflict. Kept slide with more text (global_id={best.global_id}).")
        return best
    log.append(f"Slot {slot_id}: Unknown tie '{tie}', kept first.")
    return candidates[0]


def select_slides_by_blueprint(slide_infos: List[SlideInfo], log: List[str]) -> List[Tuple[int, str, str, Optional[SlideInfo]]]:
    """
    按 SLIDE_BLUEPRINT 为每个 Slot 选出唯一优胜 Slide。
    返回: [(slot_id, chapter_id, chapter_title, slide_info), ...]，未匹配的 slot 对应 slide_info=None。
    """
    result: List[Tuple[int, str, str, Optional[SlideInfo]]] = []
    for spec in SLIDE_BLUEPRINT:
        slot_id = spec["slot"]
        ch_id = spec["chapter_id"]
        ch_title = spec["chapter_title"]
        tie = spec.get("tie", "first")
        candidates = [s for s in slide_infos if _slide_matches_blueprint_slot(s, spec)]
        winner = _tie_break(candidates, tie, slot_id, log)
        result.append((slot_id, ch_id, ch_title, winner))
    return result


def _match_slide_to_script(slide_title: str) -> Optional[Tuple[str, int, str]]:
    """
    用剧本中的标题片段匹配幻灯片：标题包含该片段即视为匹配。
    返回: (chapter_id, sub_order, matched_phrase) 或 None（未入选）。
    按剧本顺序取第一个匹配的章节与槽位（每槽位可有多条备选短语，任一匹配即可）。
    """
    if not slide_title or not slide_title.strip():
        return None
    title_norm = re.sub(r"\s+", " ", slide_title.strip())
    for ch_id, _ch_name, slots in CURATED_SCRIPT:
        for sub_order, slot_phrases in enumerate(slots):
            for phrase in slot_phrases:
                p = phrase.strip() if isinstance(phrase, str) else ""
                if not p:
                    continue
                if p in title_norm or (len(title_norm) <= len(p) and title_norm in p):
                    return (ch_id, sub_order, phrase)
    return None


def select_slides_by_curated_script(
    slide_infos: List[SlideInfo],
) -> List[SlideInfo]:
    """
    按剧本精选：每页用标题（或已挂的 _match_result，含正文）匹配剧本；同槽多页保留 shape_count 最大的一页。
    返回按 (章节顺序, 剧本内顺序) 排好序的 SlideInfo 列表。
    """
    # (chapter_id, sub_order) -> [(slide_info, matched_phrase), ...]
    slot_candidates: Dict[Tuple[str, int], List[Tuple[SlideInfo, str]]] = {}
    for s in slide_infos:
        match = getattr(s, "_match_result", None)
        if match is None:
            match = _match_slide_to_script(getattr(s, "slide_title", "") or "")
        if match is None:
            continue
        ch_id, sub_order, phrase = match
        key = (ch_id, sub_order)
        slot_candidates.setdefault(key, []).append((s, phrase))
    # 每个 slot 只保留 shape_count 最大的一页
    chapter_order = {ch_id: i for i, (ch_id, _, _) in enumerate(CURATED_SCRIPT)}
    selected: List[Tuple[int, int, SlideInfo]] = []
    for (ch_id, sub_order), candidates in slot_candidates.items():
        best = max(candidates, key=lambda x: x[0].shape_count)
        selected.append((chapter_order[ch_id], sub_order, best[0]))
    selected.sort(key=lambda x: (x[0], x[1]))
    return [s for _co, _so, s in selected]


def add_section_divider_slide(
    prs: Presentation,
    part_label: str,
    chapter_title: str,
    layout_index: Optional[int] = None,
) -> None:
    """
    在 prs 末尾插入一张过渡页：深绿背景，居中大标题 "Part X: Chapter Title"。
    使用空白版式，清空占位符、设背景、加文本框。
    """
    layouts = prs.slide_layouts
    idx = min(6, len(layouts) - 1) if layout_index is None else min(layout_index, len(layouts) - 1)
    slide_layout = layouts[idx]
    slide = prs.slides.add_slide(slide_layout)
    # 移除原有占位符
    for shape in list(slide.shapes):
        try:
            sp = shape.element
            sp.getparent().remove(sp)
        except Exception:
            pass
    # 深绿色背景（TVAX 品牌色）
    try:
        slide.follow_master_background = False
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(45, 90, 39)
    except Exception:
        pass
    # 居中大标题
    title_text = f"Part {part_label}: {chapter_title}"
    left, top, width, height = Inches(0.5), Inches(2.2), Inches(9), Inches(1.8)
    try:
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.name = "Arial"
    except Exception:
        pass


def run_editor() -> None:
    """
    「编辑」逻辑：严格按 SLIDE_BLUEPRINT 去重与排序。
    指纹扫描 -> 每 Slot 匹配 -> Tie-Breaker 选唯一胜者 -> 按 Slot 1..10 重组并插入过渡页。
    """
    input_dir = get_merge_dir()
    output_dir = get_output_dir()
    print("=" * 80)
    print("PPT 临床报告重组与标准化处理")
    print("=" * 80)
    print(f"输入目录（源 PPT）: {input_dir}")
    print(f"输出目录: {output_dir}")
    print()

    try:
        template_path, source_paths = auto_discover_ppts(input_dir)
        print(f"模板文件: {os.path.basename(template_path)}")
        print(f"源文件数量: {len(source_paths)}")
        for i, p in enumerate(source_paths, 1):
            print(f"  [{i}] {os.path.basename(p)}")
        print()

        slide_infos = build_slide_infos(source_paths)
        if not slide_infos:
            print("❌ 错误: 未发现任何幻灯片。")
            return

        print(f"✓ 成功加载 {len(slide_infos)} 张幻灯片")
        print()


        for s in slide_infos:
            s.slide_title = _get_slide_title(s.slide_obj)
            s.has_table = _slide_has_table(s.slide_obj)

        log: List[str] = []
        slot_results = select_slides_by_blueprint(slide_infos, log)
        
        # 未匹配的 Slot 按序号回退；多序号时用该 Slot 的 tie_breaker 选一
        n = len(slide_infos)
        spec_by_slot = {sp["slot"]: sp for sp in SLIDE_BLUEPRINT}
        new_results: List[Tuple[int, str, str, Optional[SlideInfo]]] = []
        for slot_id, ch_id, ch_title, winner in slot_results:
            if winner is None and slot_id in BLUEPRINT_FALLBACK_SLIDES:
                candidates = [slide_infos[num - 1] for num in BLUEPRINT_FALLBACK_SLIDES[slot_id] if 1 <= num <= n]
                if candidates:
                    spec = spec_by_slot.get(slot_id, {})
                    tie = spec.get("tie", "first")
                    winner = _tie_break(candidates, tie, slot_id, log)
                    if winner:
                        log.append(f"Slot {slot_id}: No match, fallback to Slide(s) {BLUEPRINT_FALLBACK_SLIDES[slot_id]} -> kept global_id={winner.global_id}.")
            new_results.append((slot_id, ch_id, ch_title, winner))
        slot_results = new_results

        # 构建章节映射和移动日志
        chapter_moves: Dict[str, List[Tuple[int, int, str]]] = {}  # chapter_id -> [(original_idx, new_idx, title)]
        final_slide_order: List[Tuple[int, str, Optional[SlideInfo]]] = []
        
        template_prs = Presentation(template_path)
        last_chapter: Optional[str] = None
        slide_count = 0
        current_final_index = 0

        print("=" * 80)
        print("处理日志")
        print("=" * 80)
        for line in log:
            print(f"  {line}")
        print()

        print("=" * 80)
        print("章节重组详情")
        print("=" * 80)
        
        for slot_id, ch_id, ch_title, winner in slot_results:
            if winner is None:
                continue
            
            # 记录章节切换
            if ch_id != last_chapter:
                if last_chapter is not None:
                    print()
                print(f"\nChapter {ch_id}: {ch_title}")
                add_section_divider_slide(template_prs, ch_id, ch_title)
                slide_count += 1
                current_final_index += 1
                last_chapter = ch_id
            
            # 记录幻灯片移动
            original_idx = winner.source_index + 1  # 1-based for display
            slide_title = getattr(winner, "slide_title", "") or f"Slide {winner.global_id}"
            final_slide_order.append((current_final_index, ch_id, winner))
            
            if ch_id not in chapter_moves:
                chapter_moves[ch_id] = []
            chapter_moves[ch_id].append((original_idx, current_final_index + 1, slide_title))
            
            print(f"  ✓ Kept Slide ID {winner.global_id} (Source: {winner.source_ppt}, Original Index: {original_idx}, Final Index: {current_final_index + 1})")
            print(f"    标题: {slide_title[:60]}{'...' if len(slide_title) > 60 else ''}")
            
            new_slide = clone_slide_into_presentation(template_prs, winner.slide_obj, layout_index=0)
            standardize_fonts(new_slide, title_font_name="Arial", body_font_name="Arial",
                              title_size_pt=24, body_size_pt=14)
            slide_count += 1
            current_final_index += 1

        print()
        print("=" * 80)
        print("章节移动摘要")
        print("=" * 80)
        
        chapter_names = {ch_id: ch_title for _, ch_id, ch_title, _ in slot_results}
        for ch_id in sorted(set(ch_id for _, ch_id, _, _ in slot_results if ch_id)):
            moves = chapter_moves.get(ch_id, [])
            if moves:
                ch_name = chapter_names.get(ch_id, f"Chapter {ch_id}")
                print(f"\nChapter {ch_id}: {ch_name}")
                for orig_idx, final_idx, title in moves:
                    if orig_idx != final_idx:
                        print(f"  → Moved Slide '{title[:50]}{'...' if len(title) > 50 else ''}' from Index {orig_idx} to Index {final_idx}")
                    else:
                        print(f"  ✓ Kept Slide '{title[:50]}{'...' if len(title) > 50 else ''}' at Index {orig_idx}")

        output_pptx = os.path.join(output_dir, "merged_presentation.pptx")
        template_prs.save(output_pptx)
        
        print()
        print("=" * 80)
        print("处理完成")
        print("=" * 80)
        print(f"✓ 共生成 {slide_count} 张幻灯片（含章节分隔页）")
        print(f"✓ 成品 PPT: {output_pptx}")
        
        # 自检
        _self_check(output_pptx, slide_count, len(slide_infos))
        
    except Exception as e:
        print(f"\n❌ 错误: {str(e)}")
        import traceback
        traceback.print_exc()
        raise


def run_engine() -> None:
    """
    主流程：发现源 PPT -> 提取标题 -> 按剧本精选 -> 按章节+顺序重排 -> 插入过渡页 -> 输出。
    """
    input_dir = get_merge_dir()
    output_dir = get_output_dir()
    print("=" * 80)
    print("PPT 临床报告叙事重组处理")
    print("=" * 80)
    print(f"输入目录（源 PPT）: {input_dir}")
    print(f"输出目录: {output_dir}")
    print()

    try:
        template_path, source_paths = auto_discover_ppts(input_dir)
        print(f"模板文件: {os.path.basename(template_path)}")
        print(f"源文件数量: {len(source_paths)}")
        for i, p in enumerate(source_paths, 1):
            print(f"  [{i}] {os.path.basename(p)}")
        print()

        slide_infos = build_slide_infos(source_paths)
        if not slide_infos:
            print("❌ 错误: 未发现任何幻灯片。")
            return

        print(f"✓ 成功加载 {len(slide_infos)} 张幻灯片")
        print()

        # 记录原始索引映射
        original_index_map = {s.global_id: (s.source_ppt, s.source_index + 1) for s in slide_infos}

        for s in slide_infos:
            s.slide_title = _get_slide_title(s.slide_obj)

        # 匹配时用「标题 + 正文」提高命中率（如正文含「400」「非临床剂量」等）
        def _match(s: SlideInfo) -> Optional[Tuple[str, int, str]]:
            text = (getattr(s, "slide_title", "") or "") + "\n" + (s.text_content or "")
            return _match_slide_to_script(text)

        # 临时挂到 slide 上供 select 使用
        for s in slide_infos:
            s._match_result = _match(s)
        ordered = select_slides_by_curated_script(slide_infos)
        if not ordered and CURATED_SLIDE_NUMBERS:
            # 回退：按「Slide 序号」(1-based 原始合并顺序) 从剧本取页
            ordered = []
            for ch_id, nums in CURATED_SLIDE_NUMBERS:
                for num in nums:
                    if 1 <= num <= len(slide_infos):
                        s = slide_infos[num - 1]
                        s.chapter_id = ch_id
                        s.script_phrase = f"Slide {num}"
                        ordered.append(s)
            print("⚠ 已按 CURATED_SLIDE_NUMBERS（Slide 序号）精选，未使用标题匹配。")
        if not ordered:
            print("❌ 剧本匹配结果为空，请检查 CURATED_SCRIPT 标题片段或 CURATED_SLIDE_NUMBERS 序号。")
            print("--- 前 10 页标题/正文摘要（供对照）---")
            for i, s in enumerate(slide_infos[:10]):
                tit = getattr(s, "slide_title", "") or ""
                body = (s.text_content or "")[:80].replace("\n", " ")
                print(f"  [{i+1}] 标题: {tit!r}")
                print(f"       正文: {body!r}...")
            return

        # 为精选结果标注章节（用于过渡页与报告）
        chapter_order_map = {ch_id: (i, ch_title) for i, (ch_id, ch_title, _) in enumerate(CURATED_SCRIPT)}
        for s in ordered:
            m = getattr(s, "_match_result", None)
            if m:
                s.chapter_id = m[0]
                s.script_phrase = m[2]
            else:
                s.chapter_id = ""
                s.script_phrase = ""

        # 构建章节移动日志
        chapter_moves: Dict[str, List[Tuple[int, int, str]]] = {}
        
        template_prs = Presentation(template_path)
        last_chapter: Optional[str] = None
        slide_count = 0
        current_final_index = 0

        print("=" * 80)
        print("章节重组详情")
        print("=" * 80)

        for s in ordered:
            ch_id = getattr(s, "chapter_id", "")
            if ch_id and ch_id != last_chapter:
                if last_chapter is not None:
                    print()
                _, ch_title = chapter_order_map.get(ch_id, (0, ch_id))
                print(f"\nChapter {ch_id}: {ch_title}")
                add_section_divider_slide(template_prs, ch_id, ch_title)
                slide_count += 1
                current_final_index += 1
                last_chapter = ch_id
            
            # 记录移动
            orig_ppt, orig_idx = original_index_map.get(s.global_id, ("", 0))
            slide_title = getattr(s, "slide_title", "") or f"Slide {s.global_id}"
            final_idx = current_final_index + 1
            
            if ch_id not in chapter_moves:
                chapter_moves[ch_id] = []
            chapter_moves[ch_id].append((orig_idx, final_idx, slide_title))
            
            print(f"  ✓ Kept Slide ID {s.global_id} (Source: {orig_ppt}, Original Index: {orig_idx}, Final Index: {final_idx})")
            print(f"    标题: {slide_title[:60]}{'...' if len(slide_title) > 60 else ''}")
            
            new_slide = clone_slide_into_presentation(template_prs, s.slide_obj, layout_index=0)
            standardize_fonts(new_slide, title_font_name="Arial", body_font_name="Arial",
                              title_size_pt=24, body_size_pt=14)
            slide_count += 1
            current_final_index += 1

        print()
        print("=" * 80)
        print("章节移动摘要")
        print("=" * 80)
        
        for ch_id in sorted(set(getattr(s, "chapter_id", "") for s in ordered if getattr(s, "chapter_id", ""))):
            moves = chapter_moves.get(ch_id, [])
            if moves:
                _, ch_title = chapter_order_map.get(ch_id, (0, f"Chapter {ch_id}"))
                print(f"\nChapter {ch_id}: {ch_title}")
                for orig_idx, final_idx, title in moves:
                    if orig_idx != final_idx:
                        print(f"  → Moved Slide '{title[:50]}{'...' if len(title) > 50 else ''}' from Index {orig_idx} to Index {final_idx}")
                    else:
                        print(f"  ✓ Kept Slide '{title[:50]}{'...' if len(title) > 50 else ''}' at Index {orig_idx}")

        output_pptx = os.path.join(output_dir, "merged_presentation.pptx")
        template_prs.save(output_pptx)

        # 导出叙事报告：所有源页 + 是否入选 + 剧本条目
        try:
            import pandas as pd
            rows = []
            for s in slide_infos:
                in_script = "是" if s in ordered else "否"
                rows.append({
                    "global_id": s.global_id,
                    "source_ppt": s.source_ppt,
                    "slide_title": getattr(s, "slide_title", ""),
                    "chapter_id": getattr(s, "chapter_id", ""),
                    "script_phrase": getattr(s, "script_phrase", ""),
                    "shape_count": s.shape_count,
                    "in_final_deck": in_script,
                })
            df = pd.DataFrame(rows)
            report_path = os.path.join(output_dir, "narrative_report.xlsx")
            df.to_excel(report_path, index=False)
            print(f"\n✓ 叙事报告: {report_path}")
        except Exception as e:
            print(f"\n⚠ 叙事报告导出跳过: {e}")

        n_content = len(ordered)
        n_dividers = slide_count - n_content
        
        print()
        print("=" * 80)
        print("处理完成")
        print("=" * 80)
        print(f"✓ 共生成 {slide_count} 张幻灯片（含 {n_dividers} 个过渡页，{n_content} 张内容页）")
        print(f"✓ 成品 PPT: {output_pptx}")
        
        # 自检
        _self_check(output_pptx, slide_count, len(slide_infos))
        
    except Exception as e:
        print(f"\n❌ 错误: {str(e)}")
        import traceback
        traceback.print_exc()
        raise


def _self_check(output_pptx: str, expected_slides: int, original_count: int) -> None:
    """自检功能：验证输出文件完整性"""
    print()
    print("=" * 80)
    print("自检报告")
    print("=" * 80)
    
    checks_passed = 0
    checks_total = 0
    
    # 检查1: 文件是否存在
    checks_total += 1
    if os.path.exists(output_pptx):
        file_size = os.path.getsize(output_pptx)
        print(f"✓ [检查1] 输出文件存在 ({file_size:,} 字节)")
        checks_passed += 1
    else:
        print("❌ [检查1] 输出文件不存在")
    
    # 检查2: 文件可读性
    checks_total += 1
    try:
        prs = Presentation(output_pptx)
        actual_slides = len(prs.slides)
        print("✓ [检查2] PPT文件可正常读取")
        checks_passed += 1
        
        # 检查3: 幻灯片数量
        checks_total += 1
        if actual_slides == expected_slides:
            print(f"✓ [检查3] 幻灯片数量正确 ({actual_slides} 张)")
            checks_passed += 1
        else:
            print(f"⚠ [检查3] 幻灯片数量不匹配: 期望 {expected_slides}，实际 {actual_slides}")
        
        # 检查4: 幻灯片完整性
        checks_total += 1
        all_valid = True
        for i, slide in enumerate(prs.slides, 1):
            try:
                # 检查是否有形状
                if len(slide.shapes) == 0:
                    print(f"  ⚠ 幻灯片 {i} 无形状")
                    all_valid = False
            except Exception as e:
                print(f"  ❌ 幻灯片 {i} 读取错误: {e}")
                all_valid = False
        
        if all_valid:
            print("✓ [检查4] 所有幻灯片结构完整")
            checks_passed += 1
        else:
            print("⚠ [检查4] 部分幻灯片可能存在问题")
        
    except Exception as e:
        print(f"❌ [检查2] PPT文件读取失败: {e}")
    
    # 检查5: 去重效果
    checks_total += 1
    reduction = original_count - expected_slides
    reduction_pct = (reduction / original_count * 100) if original_count > 0 else 0
    print(f"✓ [检查5] 去重统计: 原始 {original_count} 张 → 最终 {expected_slides} 张 (减少 {reduction} 张, {reduction_pct:.1f}%)")
    checks_passed += 1
    
    print()
    print(f"自检结果: {checks_passed}/{checks_total} 项检查通过")
    if checks_passed == checks_total:
        print("✓ 所有检查通过，文件处理成功！")
    else:
        print("⚠ 部分检查未通过，请检查输出文件")


if __name__ == "__main__":
    run_editor()  # 严格按 SLIDE_BLUEPRINT 去重与叙事顺序；可改为 run_engine() 使用剧本+序号回退
