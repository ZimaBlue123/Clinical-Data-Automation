#!/usr/bin/env python
"""
高保真 PPT 合并与整合脚本（03_PPT_Merge 子项目）

目录约定：
- input/：存放待合并的源 .pptx；若存在 template.pptx 则作为模板，否则按文件名取第一份为模板。
- output/：产出 merged_presentation.pptx、merge_report.xlsx。

设计说明：
1. 在 input/ 下自动发现所有 .pptx 作为候选源文件。
   - 若存在 template.pptx，则优先作为 Master 模板，其余 .pptx 为源文件；
   - 若不存在 template.pptx，则使用按文件名排序后的第一份 PPT 作为模板，其余作为源文件。
2. 使用 TF-IDF + 余弦相似度构建“内容指纹”，结合形状数量（shape_count）度量版式复杂度：
   - 相似度 < 0.6：视为独立内容，直接保留（decision = keep_unique）；
   - 0.6 <= 相似度 < 0.9：视为“同一主题不同视角”，两页都保留（decision = keep_partial_overlap），
     并在 Notes 中标记“需人工确认整合”，不做自动删减；
   - 相似度 >= 0.9：视为“高度相似冲突”，只保留形状数更多（更复杂版式）的页面（保留者正常输出，
     被舍弃者 decision = drop_conflict，不复制到成品 PPT）。
3. “保留高复杂度版式”的核心实现：
   - 对于被保留的幻灯片，使用底层 XML 深拷贝方式复制原始 slide.shapes：
       * 创建一张基于 Master 模板的空白幻灯片（继承主题/背景）；
       * 删除该新幻灯片自带的占位符；
       * 对原 slide 中的每一个 shape，将其 XML element 深拷贝后插入新幻灯片的 shape tree。
     这样可以：
       * 维持复杂布局与坐标；
       * Group Shapes（组合图形）整体复制，不打散内部元素；
       * Chart 对象（图表）仍然是可编辑对象，而不是转为图片。
4. 字体统一：
   - 标题（通过 slide.shapes.title 近似识别）字体设为 Arial/黑体，字号 24pt；
   - 其他正文文本统一为 Arial/黑体，字号 14pt。
5. 输出：
   - merged_presentation.pptx：合并后的 PPT；
   - merge_report.xlsx：每张原始幻灯片的来源、相似度与处理结果。
"""

import os
from copy import deepcopy
from dataclasses import dataclass, field
from typing import Any

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.opc.package import Part
from pptx.util import Pt
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# 用于在 XML 中识别 r:embed / r:id 的命名空间
_NS_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


# =========================
# 数据结构
# =========================


@dataclass
class SlideInfo:
    """用于保存单个幻灯片的结构化信息"""

    global_id: int  # 全局 ID（跨所有文件连续编号）
    source_ppt: str  # 源文件名
    source_path: str  # 源文件完整路径
    source_index: int  # 源文件中的 slide 序号（从 0 开始）
    text_content: str  # 文本内容指纹的原始语料
    shape_count: int  # 形状数量，用作版式复杂度指标
    slide_obj: object  # python-pptx 的 slide 对象
    cluster_id: int | None = None  # 连通分量 ID
    max_similarity: float = 0.0  # 与任意其他页的最大相似度
    decision: str = "undecided"  # keep_unique / keep_partial_overlap / drop_conflict
    notes_flags: list[str] = field(default_factory=list)  # 备注中需要写入的标记


# =========================
# 文本与形状提取
# =========================


def extract_text_from_shape(shape) -> str:
    """从单个 shape 中提取尽可能多的文本信息（文本框、表格、图表标题、alt text 等）。"""
    texts = []

    # 通用 alt_text
    try:
        if hasattr(shape, "alternative_text") and shape.alternative_text:
            texts.append(str(shape.alternative_text))
    except Exception:
        pass

    # 文本框 / 占位文本
    try:
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                texts.append(paragraph.text)
    except Exception:
        pass

    # 表格
    try:
        if getattr(shape, "has_table", False):
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    texts.append(cell.text)
    except Exception:
        pass

    # 图表标题
    try:
        if hasattr(shape, "chart"):
            chart = shape.chart
            if chart.has_title:
                texts.append(chart.chart_title.text_frame.text)
    except Exception:
        pass

    # 组合图形（GroupShape）不打散，仅复制时整体拷贝，这里不做递归文本提取

    return "\n".join(t for t in texts if t)


def extract_slide_text(slide) -> str:
    """从整张 slide 中提取文本（标题、正文、表格、图表标题、备注等）。"""
    parts = []

    # 标题占位符
    try:
        if slide.shapes.title:
            parts.append(slide.shapes.title.text)
    except Exception:
        pass

    # 其他 shape
    for shape in slide.shapes:
        try:
            parts.append(extract_text_from_shape(shape))
        except Exception:
            continue

    # Notes
    try:
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                parts.append(notes_slide.notes_text_frame.text)
    except Exception:
        pass

    return "\n".join(p for p in parts if p)


def count_shapes(slide) -> int:
    """统计 slide 中 shape 的数量，作为版式复杂度近似指标。"""
    try:
        return len(slide.shapes)
    except Exception:
        return 0


# =========================
# 幻灯片复制与字体标准化（含图片/图表 part 与 rId 重映射）
# =========================


def _collect_rids_from_element(element) -> set[str]:
    """从 XML 元素树中收集所有 r:embed / r:id 属性值（关系 ID）。"""
    rids: set[str] = set()
    for el in element.iter():
        for key, value in list(el.attrib.items()):
            if "}" in key:
                ns, local = key.split("}", 1)
                ns = ns.strip("{")
                if ns == _NS_REL and local in ("embed", "id") and value:
                    rids.add(value)
            elif key in ("embed", "id") and value:
                rids.add(value)
    return rids


def _replace_rids_in_element(element, rid_map: dict[str, str]) -> None:
    """将元素树中所有 r:embed / r:id 属性值按 rid_map 替换为新的 rId（原地修改）。"""
    for el in element.iter():
        for key in list(el.attrib.keys()):
            if "}" in key:
                ns, local = key.split("}", 1)
                ns = ns.strip("{")
                if ns == _NS_REL and local in ("embed", "id"):
                    old = el.attrib[key]
                    if old in rid_map:
                        el.attrib[key] = rid_map[old]
            elif key in ("embed", "id") and el.attrib[key] in rid_map:
                el.attrib[key] = rid_map[el.attrib[key]]


def _copy_image_part_to_package(source_part: Part, target_package) -> Part | None:
    """将源包中的图片 part 复制到目标包，返回新 part（供 relate_to 使用）。"""
    try:
        blob = source_part.blob
        ext = getattr(source_part.partname, "ext", None) or "png"
        if ext not in ("png", "jpg", "jpeg", "gif", "bmp", "tiff", "emf", "wmf"):
            ext = "png"
        partname = target_package.next_image_partname(ext)
        return Part(
            partname,
            source_part.content_type,
            target_package,
            blob=blob,
        )
    except Exception:
        return None


def _copy_chart_part_to_package(
    source_chart_part, source_slide_part, target_package, target_slide_part
) -> tuple[Any, str] | None:
    """
    将源图表 part（及其嵌入的 xlsx）复制到目标包。
    返回 (新 ChartPart, 新 rId) 或 None。新 chart 已与 target_slide_part 建立关系。
    """
    try:
        from pptx.opc.package import XmlPart

        blob = source_chart_part.blob
        xlsx_rId = None
        try:
            if hasattr(source_chart_part._element, "chartSpace"):
                cspace = source_chart_part._element.chartSpace
            else:
                cspace = source_chart_part._element
            if hasattr(cspace, "externalData") and cspace.externalData is not None:
                xlsx_rId = cspace.externalData.rId
        except Exception:
            pass
        if xlsx_rId is None and hasattr(source_chart_part._element, "xlsx_part_rId"):
            xlsx_rId = source_chart_part._element.xlsx_part_rId

        new_xlsx_part = None
        if xlsx_rId:
            try:
                xlsx_part = source_chart_part.related_part(xlsx_rId)
                partname_xlsx = target_package.next_partname("/ppt/embeddings/Microsoft_Excel_Sheet%d.xlsx")
                new_xlsx_part = Part(
                    partname_xlsx,
                    CT.SML_SHEET,
                    target_package,
                    blob=xlsx_part.blob,
                )
            except Exception:
                pass

        partname_chart = target_package.next_partname("/ppt/charts/chart%d.xml")
        new_chart_part = XmlPart.load(partname_chart, CT.DML_CHART, target_package, blob)
        if new_xlsx_part is not None:
            new_xlsx_rId = new_chart_part.relate_to(new_xlsx_part, RT.PACKAGE)
            try:
                ed = getattr(new_chart_part._element, "externalData", None)
                if ed is not None:
                    ed.rId = new_xlsx_rId
            except Exception:
                pass
        new_rId = target_slide_part.relate_to(new_chart_part, RT.CHART)
        return (new_chart_part, new_rId)
    except Exception:
        return None


def _copy_part_to_target_and_get_new_rid(  # noqa: PLR0912 - TODO: 下个迭代重构 # noqa: PLR0911 - TODO: 下个迭代重构
    source_slide_part,
    target_slide_part,
    target_package,
    rId: str,
    rid_map: dict[str, str],
) -> str | None:
    """
    若 rId 尚未在 rid_map 中，则将 source_slide_part 通过 rId 引用的 part 复制到
    target_package，并在 target_slide_part 上建立关系，返回新 rId 并写入 rid_map。
    不复制 SLIDE_LAYOUT / NOTES_SLIDE 等与幻灯片结构相关的关系。
    """
    if rId in rid_map:
        return rid_map[rId]
    try:
        rel = source_slide_part.rels[rId]
    except KeyError:
        return None
    reltype = rel.reltype
    if reltype in (RT.SLIDE_LAYOUT, RT.NOTES_SLIDE):
        return None
    source_part = rel.target_part

    try:
        if reltype == RT.IMAGE:
            new_part = _copy_image_part_to_package(source_part, target_package)
            if new_part is not None:
                new_rId = target_slide_part.relate_to(new_part, RT.IMAGE)
                rid_map[rId] = new_rId
                return new_rId
        elif reltype == RT.CHART:
            result = _copy_chart_part_to_package(source_part, source_slide_part, target_package, target_slide_part)
            if result is not None:
                _, new_rId = result
                rid_map[rId] = new_rId
                return new_rId
        elif reltype in (RT.PACKAGE, RT.OLE_OBJECT):
            try:
                blob = source_part.blob
                ext = getattr(source_part.partname, "ext", "bin")
                if "xlsx" in getattr(source_part.partname, "filename", ""):
                    partname = target_package.next_partname("/ppt/embeddings/oleObject%d.xlsx")
                else:
                    partname = target_package.next_partname("/ppt/embeddings/oleObject%d.bin")
                new_part = Part(
                    partname,
                    source_part.content_type,
                    target_package,
                    blob=blob,
                )
                new_rId = target_slide_part.relate_to(new_part, reltype)
                rid_map[rId] = new_rId
                return new_rId
            except Exception:
                pass
        elif reltype in (RT.VIDEO, RT.MEDIA):
            try:
                blob = getattr(source_part, "blob", None)
                if blob is not None:
                    ext = getattr(source_part.partname, "ext", "bin")
                    partname = target_package.next_media_partname(ext)
                    new_part = Part(
                        partname,
                        source_part.content_type,
                        target_package,
                        blob=blob,
                    )
                    new_rId = target_slide_part.relate_to(new_part, reltype)
                    rid_map[rId] = new_rId
                    return new_rId
            except Exception:
                pass
    except Exception:
        pass
    return None


def clone_slide_into_presentation(template_prs: Presentation, source_slide, layout_index: int = 0):
    """
    将 source_slide 的内容复制到 template_prs 中的一个新幻灯片。

    “保留高复杂度版式”的关键逻辑：
    - 新建幻灯片时使用模板的某个 layout（继承母版/主题/背景）；
    - 删除 layout 自带占位符后，对源 slide 的每个 shape：深拷贝其 XML，并**将该 shape 引用的
      所有 part（图片、图表、OLE 等）复制到目标包并重映射 rId**，再插入新 slide，从而避免
      “无法显示该图片”等断链问题；
    - 不对 Group Shapes 做拆分，整个 XML 树一并复制；
    - 图表对象（Chart）及其嵌入的 xlsx 一并复制，仍为可编辑图表。
    """
    slide_layout = template_prs.slide_layouts[layout_index]
    new_slide = template_prs.slides.add_slide(slide_layout)
    target_slide_part = new_slide.part
    target_package = template_prs.part.package
    source_slide_part = source_slide.part

    # 删除 layout 自带的占位符
    for shape in list(new_slide.shapes):
        el = shape.element
        el.getparent().remove(el)

    # 对每个 shape：收集 rId -> 复制 part 并建立新关系 -> 替换 XML 中的 rId -> 插入
    for shape in source_slide.shapes:
        try:
            new_el = deepcopy(shape.element)
            rids = _collect_rids_from_element(new_el)
            rid_map: dict[str, str] = {}
            for rId in rids:
                _copy_part_to_target_and_get_new_rid(
                    source_slide_part,
                    target_slide_part,
                    target_package,
                    rId,
                    rid_map,
                )
            _replace_rids_in_element(new_el, rid_map)
            new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
        except Exception:
            continue

    # 不再复制背景（易导致主题/版式错乱）；若需可再按 part 复制
    return new_slide


def standardize_fonts(slide, title_font_name="Arial", body_font_name="Arial", title_size_pt=24, body_size_pt=14):
    """
    对复制后的幻灯片统一字体：
    - 标题：优先识别 slide.shapes.title 所在 shape，设为标题字体和较大字号；
    - 其他文本：统一为正文字体和 12–14pt 左右字号。
    """
    title_ids = set()
    try:
        if slide.shapes.title:
            title_ids.add(slide.shapes.title.shape_id)
    except Exception:
        pass

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        is_title = shape.shape_id in title_ids
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                if is_title:
                    font.name = title_font_name
                    font.size = Pt(title_size_pt)
                else:
                    font.name = body_font_name
                    font.size = Pt(body_size_pt)


def add_notes_to_slide(slide, extra_notes: list[str]):
    """在 Notes 中追加“需人工确认整合”等提示信息。"""
    if not extra_notes:
        return

    text = "需人工确认整合\n" + "\n".join(extra_notes)

    notes_slide = slide.notes_slide  # 访问属性时若不存在会自动创建
    if notes_slide.notes_text_frame.text:
        notes_slide.notes_text_frame.text += "\n\n" + text
    else:
        notes_slide.notes_text_frame.text = text


# =========================
# 相似度与聚类决策
# =========================


def build_slide_infos(ppt_paths: list[str]) -> list[SlideInfo]:
    """从多个 PPT 文件中构造 SlideInfo 列表。"""
    slide_infos: list[SlideInfo] = []
    gid = 0

    for path in ppt_paths:
        prs = Presentation(path)
        for idx, slide in enumerate(prs.slides):
            text = extract_slide_text(slide)
            shapes = count_shapes(slide)
            slide_infos.append(
                SlideInfo(
                    global_id=gid,
                    source_ppt=os.path.basename(path),
                    source_path=path,
                    source_index=idx,
                    text_content=text if text.strip() else "",
                    shape_count=shapes,
                    slide_obj=slide,
                )
            )
            gid += 1

    return slide_infos


def compute_similarity(slide_infos: list[SlideInfo]) -> np.ndarray:
    """基于 TF-IDF 计算所有幻灯片之间的余弦相似度矩阵。"""
    corpus = [s.text_content for s in slide_infos]
    if all(not c.strip() for c in corpus):
        # 文本全空时，返回单位矩阵，避免 TF-IDF 出错
        return np.eye(len(slide_infos))

    vectorizer = TfidfVectorizer(
        max_features=5000,
        ngram_range=(1, 2),
        stop_words=None,
    )
    tfidf = vectorizer.fit_transform(corpus)
    return cosine_similarity(tfidf)


def build_clusters(sim_matrix: np.ndarray, low: float = 0.6) -> dict[int, int]:
    """
    使用相似度 >= low 的边构造连通分量，得到 cluster_id。
    这里使用简单并查集，避免真正的 KMeans 需要指定簇数的问题。
    """
    n = sim_matrix.shape[0]
    parent = list(range(n))

    def find(x):
        while parent[x] != x:
            parent[x] = parent[parent[x]]
            x = parent[x]
        return x

    def union(a, b):
        ra, rb = find(a), find(b)
        if ra != rb:
            parent[rb] = ra

    for i in range(n):
        for j in range(i + 1, n):
            if sim_matrix[i, j] >= low:
                union(i, j)

    root_to_cluster: dict[int, int] = {}
    mapping: dict[int, int] = {}
    next_cid = 0
    for i in range(n):
        root = find(i)
        if root not in root_to_cluster:
            root_to_cluster[root] = next_cid
            next_cid += 1
        mapping[i] = root_to_cluster[root]

    return mapping


def decide_keep_drop(
    slide_infos: list[SlideInfo],
    sim_matrix: np.ndarray,
    thr_low: float = 0.6,
    thr_high: float = 0.9,
):
    """
    根据相似度与形状复杂度，对每张幻灯片打上决策标签：
      - drop_conflict：在 >= thr_high 的“高度相似冲突”里被判定为更简单版式的那一页；
      - keep_partial_overlap：存在 0.6~0.9 区间的配对关系；保留且在 Notes 中写入“需人工确认整合”；
      - keep_unique：与任何其他页相似度 < thr_low 的独立内容。
    """
    n = len(slide_infos)

    # 记录最大相似度（排除自身）
    for i in range(n):
        sims = np.delete(sim_matrix[i], i)
        slide_infos[i].max_similarity = float(sims.max()) if sims.size > 0 else 0.0

    # 初始默认全部 keep_unique
    for s in slide_infos:
        s.decision = "keep_unique"

    # 1. 高度相似冲突：只保留形状更多的那一页
    for i in range(n):
        for j in range(i + 1, n):
            sim_ij = sim_matrix[i, j]
            if sim_ij >= thr_high:
                si, sj = slide_infos[i], slide_infos[j]
                if si.shape_count >= sj.shape_count:
                    sj.decision = "drop_conflict"
                else:
                    si.decision = "drop_conflict"

    # 2. 部分重叠：0.6 <= sim < 0.9，全保留并写备注
    for i in range(n):
        for j in range(i + 1, n):
            sim_ij = sim_matrix[i, j]
            if thr_low <= sim_ij < thr_high:
                si, sj = slide_infos[i], slide_infos[j]
                if si.decision != "drop_conflict":
                    si.decision = "keep_partial_overlap"
                    si.notes_flags.append(f"部分重叠：与全局ID {sj.global_id} 相似度 {sim_ij:.3f}")
                if sj.decision != "drop_conflict":
                    sj.decision = "keep_partial_overlap"
                    sj.notes_flags.append(f"部分重叠：与全局ID {si.global_id} 相似度 {sim_ij:.3f}")


# =========================
# 报表生成
# =========================


def generate_report(slide_infos: list[SlideInfo], out_xlsx: str):
    """输出 merge_report.xlsx，便于后续人工审阅与追溯。"""
    rows = []
    for s in slide_infos:
        rows.append(
            {
                "global_id": s.global_id,
                "source_ppt": s.source_ppt,
                "source_path": s.source_path,
                "source_slide_index": s.source_index,
                "cluster_id": s.cluster_id,
                "shape_count": s.shape_count,
                "max_similarity": s.max_similarity,
                "decision": s.decision,
            }
        )
    df = pd.DataFrame(rows)
    df = df.sort_values(by=["cluster_id", "global_id"]).reset_index(drop=True)
    df.to_excel(out_xlsx, index=False)


# =========================
# 主流程（自动发现 PPT）
# =========================

# 合并用目录：输入和输出分离
INPUT_DIR_NAME = "input"
OUTPUT_DIR_NAME = "output"


def get_input_dir() -> str:
    """返回输入目录的绝对路径，不存在则创建。"""
    script_dir = os.path.abspath(os.path.dirname(__file__))
    input_dir = os.path.join(script_dir, INPUT_DIR_NAME)
    os.makedirs(input_dir, exist_ok=True)
    return input_dir


def get_output_dir() -> str:
    """返回输出目录的绝对路径，不存在则创建。"""
    script_dir = os.path.abspath(os.path.dirname(__file__))
    output_dir = os.path.join(script_dir, OUTPUT_DIR_NAME)
    os.makedirs(output_dir, exist_ok=True)
    return output_dir


def get_merge_dir() -> str:
    """与 get_input_dir 一致，供 ppt_engine 等调用时从 input 目录发现源 PPT。"""
    return get_input_dir()


def auto_discover_ppts(base_dir: str):
    """
    自动发现 base_dir 下的所有 .pptx。
    - 若存在 template.pptx，则用其作为模板；
    - 否则，将按名称排序后的第一份 pptx 作为模板，其余为源文件；
    - 排除最终输出文件名（若已经存在旧版本）。
    """
    all_files = [f for f in os.listdir(base_dir) if f.lower().endswith(".pptx")]

    # 排除潜在旧输出
    exclude_names = {"merged_presentation.pptx"}
    files = [f for f in all_files if f not in exclude_names]

    if not files:
        raise FileNotFoundError("当前目录未发现任何 .pptx 文件。")

    template_name = None
    if "template.pptx" in files:
        template_name = "template.pptx"
    else:
        # 没有明确模板时，使用排序后的第一份 PPT 作为 Master 模板
        files.sort()
        template_name = files[0]

    template_path = os.path.join(base_dir, template_name)

    source_paths = []
    for f in files:
        if f == template_name:
            continue
        source_paths.append(os.path.join(base_dir, f))

    if not source_paths:
        # 只有一份 PPT 时，既当模板又当源，用于自复制（理论上少见）
        source_paths = [template_path]

    return template_path, source_paths


def main():
    input_dir = get_input_dir()
    output_dir = get_output_dir()
    print(f"输入目录（源 PPT）: {input_dir}")
    print(f"输出目录（交付产物）: {output_dir}")

    # 自动发现模板与源 PPT
    template_path, source_paths = auto_discover_ppts(input_dir)
    print("使用的模板文件:", template_path)
    print("发现的源 PPT 文件:")
    for p in source_paths:
        print("  -", p)

    # 构建 slide 信息
    slide_infos = build_slide_infos(source_paths)
    if not slide_infos:
        print("未从源文件中找到任何幻灯片。")
        return

    # 计算相似度矩阵
    sim_matrix = compute_similarity(slide_infos)

    # 构建 cluster_id
    cluster_mapping = build_clusters(sim_matrix, low=0.6)
    for s in slide_infos:
        s.cluster_id = cluster_mapping[s.global_id]

    # 决策：独立/部分重叠/冲突
    decide_keep_drop(slide_infos, sim_matrix, thr_low=0.6, thr_high=0.9)

    # 载入模板演示文稿
    template_prs = Presentation(template_path)

    # 将同一 cluster 的幻灯片尽量连续排列
    sorted_infos = sorted(slide_infos, key=lambda s: (s.cluster_id, s.global_id))

    # 复制保留的幻灯片
    kept_count = 0
    for s in sorted_infos:
        if s.decision == "drop_conflict":
            continue

        new_slide = clone_slide_into_presentation(template_prs, s.slide_obj, layout_index=0)

        # 部分重叠的 slide 在 Notes 中加入人工审核提示
        if s.decision == "keep_partial_overlap" and s.notes_flags:
            add_notes_to_slide(new_slide, s.notes_flags)

        # 字体标准化（如需中文黑体可在此处替换字体名）
        standardize_fonts(
            new_slide,
            title_font_name="Arial",
            body_font_name="Arial",
            title_size_pt=24,
            body_size_pt=14,
        )

        kept_count += 1

    output_pptx = os.path.join(output_dir, "merged_presentation.pptx")
    template_prs.save(output_pptx)

    output_report = os.path.join(output_dir, "merge_report.xlsx")
    generate_report(slide_infos, output_report)

    print(f"合并完成，共写入 {kept_count} 张幻灯片。")
    print(f"成品 PPT 路径: {output_pptx}")
    print(f"合并报告路径: {output_report}")


if __name__ == "__main__":
    main()
