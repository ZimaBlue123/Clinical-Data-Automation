#!/usr/bin/env python
"""
CSR 规范 PPT 整合脚本
将多个具有内容交叉、重复且版式不一的 PPT 源文件整合为一份逻辑严密、视觉统一、符合临床研究报告（CSR）规范的终极演示文稿。

功能：
1. 叙事逻辑重组 - 按5章CSR结构分类和排序
2. 智能去重与精选 - TF-IDF相似度 + 优先级算法
3. 全局视觉标准化 - 统一标题位置、字体、章节过渡页
4. 输出最终报告和详细日志
"""

from __future__ import annotations

import os
import logging
from collections import defaultdict

from merge_ppt import (
    get_merge_dir,
    get_output_dir,
    auto_discover_ppts,
    build_slide_infos,
    clone_slide_into_presentation,
    standardize_fonts,
    SlideInfo,
    extract_slide_text,
    compute_similarity,
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 配置日志
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s", datefmt="%Y-%m-%d %H:%M:%S")
logger = logging.getLogger(__name__)


# =============================================================================
# 第一章：研究设计 (Study Design)
# =============================================================================
CHAPTER_1_KEYWORDS = {
    "设计": ["设计", "Design", "design"],
    "入组": ["入组", "Enrollment", "enrollment", "招募"],
    "流程": ["流程", "Flow", "flow", "流程图"],
    "基线": ["基线", "Baseline", "baseline"],
}

CHAPTER_1_PRIORITY = {
    "contains_n64_n32": ["N=64", "N=32", "n=64", "n=32", "64例", "32例"],
    "has_table": True,
    "has_diagram": True,
}


# =============================================================================
# 第二章：安全性总体概览 (Top-line Safety)
# =============================================================================
CHAPTER_2_KEYWORDS = {
    "总体": ["总体", "Overview", "overview", "概览"],
    "摘要": ["摘要", "Summary", "summary"],
    "核心数字": ["0", "98%", "61%", "0例", "SAE", "严重不良事件"],
    "风险": ["风险", "Risk", "risk", "可控", "可控性"],
}

CHAPTER_2_PRIORITY = {
    "contains_core_numbers": ["0", "98%", "61%", "0例"],
    "has_summary": True,
}


# =============================================================================
# 第三章：反应原性详述 (Detailed Reactogenicity)
# =============================================================================
CHAPTER_3_KEYWORDS = {
    "征集性": ["征集性", "Solicited", "solicited"],
    "局部": ["局部", "Local", "local", "局部反应"],
    "全身": ["全身", "Systemic", "systemic", "全身反应"],
    "MedDRA": ["MedDRA", "meddra", "首选术语", "PT"],
    "严重程度": ["严重程度", "Grade", "grade", "Severity", "severity", "级别", "分级"],
    "持续时间": ["持续时间", "Duration", "duration", "持续"],
}

CHAPTER_3_PRIORITY = {
    "has_table": True,  # 优先保留包含表格的页面
    "has_meddra_data": True,
    "has_grade_info": True,
}


# =============================================================================
# 第四章：竞品横向对标 (Benchmark vs Shingrix)
# =============================================================================
CHAPTER_4_KEYWORDS = {
    "Shingrix": ["Shingrix", "shingrix", "欣安立适"],
    "GSK": ["GSK", "gsk", "葛兰素史克"],
    "对比": ["对比", "Comparison", "comparison", "比较", "vs", "VS"],
    "优势": ["优势", "Advantage", "advantage", "优于", "更好"],
}

CHAPTER_4_PRIORITY = {
    "contains_shingrix": True,  # 强制规则：所有提及竞品的页面必须归集至此
    "has_comparison_table": True,
}


# =============================================================================
# 第五章：总结与结论 (Conclusion)
# =============================================================================
CHAPTER_5_KEYWORDS = {
    "总结": ["总结", "Summary", "summary", "结论"],
    "价值": ["价值", "Value", "value", "临床价值", "意义"],
    "依从性": ["依从性", "Compliance", "compliance", "耐受性"],
    "下一步": ["下一步", "Next", "next", "计划", "Plan", "plan"],
}

CHAPTER_5_PRIORITY = {
    "has_conclusion": True,
    "has_value_proposition": True,
}


# =============================================================================
# 章节定义
# =============================================================================
CHAPTERS = [
    {
        "id": "1",
        "title": "研究设计 (Study Design)",
        "keywords": CHAPTER_1_KEYWORDS,
        "priority": CHAPTER_1_PRIORITY,
    },
    {
        "id": "2",
        "title": "安全性总体概览 (Top-line Safety)",
        "keywords": CHAPTER_2_KEYWORDS,
        "priority": CHAPTER_2_PRIORITY,
    },
    {
        "id": "3",
        "title": "反应原性详述 (Detailed Reactogenicity)",
        "keywords": CHAPTER_3_KEYWORDS,
        "priority": CHAPTER_3_PRIORITY,
    },
    {
        "id": "4",
        "title": "竞品横向对标 (Benchmark vs Shingrix)",
        "keywords": CHAPTER_4_KEYWORDS,
        "priority": CHAPTER_4_PRIORITY,
    },
    {
        "id": "5",
        "title": "总结与结论 (Conclusion)",
        "keywords": CHAPTER_5_KEYWORDS,
        "priority": CHAPTER_5_PRIORITY,
    },
]


# =============================================================================
# 辅助函数
# =============================================================================


def _slide_has_table(slide) -> bool:
    """判断幻灯片是否包含表格。"""
    try:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                return True
    except Exception:
        pass
    return False


def _get_slide_title(slide) -> str:
    """从单页 Slide 提取标题。"""
    try:
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip()
    except Exception:
        pass
    full = extract_slide_text(slide)
    if not full:
        return ""
    first_line = full.split("\n")[0].strip()
    return first_line[:200] if first_line else ""


def _slide_contains_keywords(slide_info: SlideInfo, keywords: dict[str, list[str]]) -> bool:
    """检查幻灯片是否包含关键词。"""
    title = (getattr(slide_info, "slide_title", "") or "").lower()
    text = (slide_info.text_content or "").lower()
    combined = title + " " + text

    # 如果文本为空，返回False
    if not combined.strip():
        return False

    # 检查是否包含任何关键词
    for _key, keyword_list in keywords.items():
        for keyword in keyword_list:
            keyword_lower = keyword.lower()
            # 更灵活的匹配：包含关键词或关键词的一部分
            if keyword_lower in combined or combined in keyword_lower:
                return True
            # 也检查部分匹配（至少3个字符）
            if len(keyword_lower) >= 3:
                for word in combined.split():
                    if keyword_lower[:3] in word or word in keyword_lower[:3]:
                        return True
    return False


def _calculate_slide_score(slide_info: SlideInfo, priority: dict) -> int:  # noqa: PLR0912 - TODO: 下个迭代重构
    """计算幻灯片优先级得分。"""
    score = 0

    # 基础分：形状数量（信息量）
    score += slide_info.shape_count * 10

    # 包含表格
    if priority.get("has_table", False) and getattr(slide_info, "has_table", False):
        score += 100

    # 包含特定数字
    if "contains_n64_n32" in priority:
        text = (slide_info.text_content or "").lower()
        for num in priority["contains_n64_n32"]:
            if num.lower() in text:
                score += 50
                break

    # 包含核心数字
    if "contains_core_numbers" in priority:
        text = (slide_info.text_content or "").lower()
        for num in priority["contains_core_numbers"]:
            if num.lower() in text:
                score += 50
                break

    # MedDRA数据
    if priority.get("has_meddra_data", False):
        text = (slide_info.text_content or "").lower()
        if "meddra" in text or "首选术语" in text:
            score += 50

    # 包含Shingrix（强制规则）
    if priority.get("contains_shingrix", False):
        text = (slide_info.text_content or "").lower()
        if "shingrix" in text or "欣安立适" in text:
            score += 200  # 高分确保归集

    # 对比表格
    if priority.get("has_comparison_table", False) and getattr(slide_info, "has_table", False):
        text = (slide_info.text_content or "").lower()
        if "对比" in text or "comparison" in text:
            score += 100

    return score


def classify_slide_to_chapter(slide_info: SlideInfo) -> dict | None:
    """将幻灯片分类到对应章节。"""
    # 获取标题和文本内容
    title = (getattr(slide_info, "slide_title", "") or "").lower()
    text = (slide_info.text_content or "").lower()
    combined = title + " " + text

    # 如果文本为空，返回None
    if not combined.strip():
        return None

    # 强制规则：第四章优先（所有提及Shingrix的必须归集）
    if _slide_contains_keywords(slide_info, CHAPTER_4_KEYWORDS):
        return CHAPTERS[3]  # 第四章

    # 按顺序检查其他章节
    for chapter in CHAPTERS:
        if chapter["id"] == "4":  # 跳过第四章（已检查）
            continue
        if _slide_contains_keywords(slide_info, chapter["keywords"]):
            return chapter

    # 如果无法分类，尝试基于文本内容的启发式分类
    # 检查是否包含数字（可能是数据页）
    if any(char.isdigit() for char in combined):
        # 包含"安全"、"不良事件"等 -> 第三章
        if any(kw in combined for kw in ["安全", "safety", "不良", "adverse", "ae", "反应"]):
            return CHAPTERS[2]  # 第三章
        # 包含百分比、数字 -> 第二章或第三章
        if any(kw in combined for kw in ["%", "percent", "百分比", "总体", "overview"]):
            return CHAPTERS[1]  # 第二章

    return None


def deduplicate_slides(slide_infos: list[SlideInfo], similarity_threshold: float = 0.85) -> list[SlideInfo]:
    """智能去重：使用TF-IDF相似度算法。"""
    if len(slide_infos) <= 1:
        return slide_infos

    # 计算相似度矩阵
    sim_matrix = compute_similarity(slide_infos)

    # 标记要删除的幻灯片
    to_remove: set[int] = set()

    for i in range(len(slide_infos)):
        if i in to_remove:
            continue

        for j in range(i + 1, len(slide_infos)):
            if j in to_remove:
                continue

            similarity = sim_matrix[i][j]
            if similarity >= similarity_threshold:
                # 相似度很高，保留得分更高的
                score_i = _calculate_slide_score(slide_infos[i], {})
                score_j = _calculate_slide_score(slide_infos[j], {})

                if score_i >= score_j:
                    to_remove.add(j)
                    logger.info(f"Deduplication: Remove slide {j} (similarity {similarity:.3f} with {i}, lower score)")
                else:
                    to_remove.add(i)
                    logger.info(f"Deduplication: Remove slide {i} (similarity {similarity:.3f} with {j}, lower score)")
                    break

    # 返回保留的幻灯片
    kept = [slide_infos[i] for i in range(len(slide_infos)) if i not in to_remove]
    logger.info(
        f"Deduplication complete: Original {len(slide_infos)} slides, kept {len(kept)}, removed {len(to_remove)}"
    )  # noqa: E501
    return kept


def standardize_slide_layout(slide) -> None:
    """标准化幻灯片布局：统一标题位置和字体。"""
    try:
        # 统一标题位置（左上角 Left=0.5", Top=0.3"）
        title_left = Inches(0.5)
        title_top = Inches(0.3)

        # 查找标题形状
        title_shape = None
        for shape in slide.shapes:
            try:
                if shape == slide.shapes.title:
                    title_shape = shape
                    break
            except Exception:
                pass

        # 如果找到标题，调整位置
        if title_shape:
            try:
                title_shape.left = title_left
                title_shape.top = title_top
            except Exception:
                pass

        # 标准化字体
        standardize_fonts(
            slide,
            title_font_name="Arial",
            body_font_name="Arial",
            title_size_pt=24,
            body_size_pt=14,
        )

    except Exception as e:
        logger.warning(f"Error standardizing layout: {e}")


def add_chapter_divider_slide(prs: Presentation, chapter_id: str, chapter_title: str) -> None:
    """添加章节过渡页：深蓝色背景，居中显示章节名称。"""
    try:
        # 使用空白版式
        blank_layout = prs.slide_layouts[6]  # 空白版式
        slide = prs.slides.add_slide(blank_layout)

        # 移除原有占位符
        for shape in list(slide.shapes):
            try:
                sp = shape.element
                sp.getparent().remove(sp)
            except Exception:
                pass

        # 深蓝色背景
        try:
            slide.follow_master_background = False
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(30, 58, 138)  # 深蓝色 #1E3A8A
        except Exception:
            pass

        # 居中大标题
        title_text = f"PART {chapter_id}: {chapter_title}"
        left = Inches(0.5)
        top = Inches(2.5)
        width = Inches(9)
        height = Inches(2)

        try:
            box = slide.shapes.add_textbox(left, top, width, height)
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
            p.font.size = Pt(36)
            p.font.name = "Arial"
            p.font.color.rgb = RGBColor(255, 255, 255)  # 白色文字
        except Exception as e:
            logger.warning(f"Error adding chapter divider: {e}")

    except Exception as e:
        logger.error(f"Failed to create chapter divider: {e}")


def integrate_ppts() -> None:  # noqa: PLR0915 - TODO: 下个迭代重构 # noqa: PLR0912 - TODO: 下个迭代重构
    """主函数：整合PPT文件。"""
    input_dir = get_merge_dir()
    output_dir = get_output_dir()

    print("=" * 80)
    print("CSR 规范 PPT 整合处理")
    print("=" * 80)
    print(f"输入目录: {input_dir}")
    print(f"输出目录: {output_dir}")
    print()

    # 自动发现PPT文件
    try:
        template_path, source_paths = auto_discover_ppts(input_dir)
        print(f"模板文件: {os.path.basename(template_path)}")
        print(f"源文件数量: {len(source_paths)}")
        for i, p in enumerate(source_paths, 1):
            print(f"  [{i}] {os.path.basename(p)}")
        print()
    except Exception as e:
        logger.error(f"发现PPT文件失败: {e}")
        return

    # 构建幻灯片信息
    try:
        slide_infos = build_slide_infos(source_paths)
        if not slide_infos:
            print("[ERROR] No slides found.")
            return
        print(f"[OK] Successfully loaded {len(slide_infos)} slides")
        print()
    except Exception as e:
        logger.error(f"构建幻灯片信息失败: {e}")
        return

    # 提取标题和表格信息
    for s in slide_infos:
        s.slide_title = _get_slide_title(s.slide_obj)
        s.has_table = _slide_has_table(s.slide_obj)

    # 智能去重
    print("=" * 80)
    print("步骤 1: 智能去重")
    print("=" * 80)
    deduplicated_slides = deduplicate_slides(slide_infos, similarity_threshold=0.85)
    print(f"去重后剩余: {len(deduplicated_slides)} 张幻灯片\n")

    # 按章节分类
    print("=" * 80)
    print("步骤 2: 按章节分类")
    print("=" * 80)
    chapter_slides: dict[str, list[SlideInfo]] = defaultdict(list)
    unclassified: list[SlideInfo] = []

    # 先显示所有幻灯片标题用于调试
    print("\n所有幻灯片标题:")
    for i, slide_info in enumerate(deduplicated_slides, 1):
        title = slide_info.slide_title[:60] if slide_info.slide_title else "[无标题]"
        print(f"  [{i}] {title}")
    print()

    for slide_info in deduplicated_slides:
        chapter = classify_slide_to_chapter(slide_info)
        if chapter:
            chapter_slides[chapter["id"]].append(slide_info)
            title = slide_info.slide_title[:50] if slide_info.slide_title else "[无标题]"
            print(f"  章节 {chapter['id']}: {title}...")
        else:
            unclassified.append(slide_info)
            title = slide_info.slide_title[:50] if slide_info.slide_title else "[无标题]"
            print(f"  [未分类]: {title}...")

    print("\n分类结果:")
    for chapter in CHAPTERS:
        count = len(chapter_slides.get(chapter["id"], []))
        print(f"  章节 {chapter['id']} ({chapter['title']}): {count} 张")
    print(f"  未分类: {len(unclassified)} 张\n")

    # 如果未分类的幻灯片太多，将它们分配到默认章节
    # 优先分配到第三章（反应原性详述），因为通常数据页最多
    if len(unclassified) > 0:
        print(f"注意: {len(unclassified)} 张幻灯片未明确分类")
        # 尝试基于内容进一步分类
        for slide_info in unclassified[:]:  # 使用切片复制列表
            text = (slide_info.text_content or "").lower()
            title = (getattr(slide_info, "slide_title", "") or "").lower()
            combined = title + " " + text

            # 启发式分类
            if "shingrix" in combined or "gsk" in combined or "对比" in combined:
                chapter_slides["4"].append(slide_info)
                unclassified.remove(slide_info)
            elif any(kw in combined for kw in ["设计", "design", "入组", "enrollment", "流程", "flow"]):
                chapter_slides["1"].append(slide_info)
                unclassified.remove(slide_info)
            elif any(kw in combined for kw in ["总结", "summary", "结论", "conclusion", "价值", "value"]):
                chapter_slides["5"].append(slide_info)
                unclassified.remove(slide_info)
            elif any(kw in combined for kw in ["总体", "overview", "摘要", "summary", "0例", "sae"]):
                chapter_slides["2"].append(slide_info)
                unclassified.remove(slide_info)
            else:
                # 默认分配到第三章（通常数据页最多）
                chapter_slides["3"].append(slide_info)
                unclassified.remove(slide_info)

        print("启发式分类后，所有幻灯片已分配到章节")

    # 每章节内按优先级排序和精选
    print("=" * 80)
    print("步骤 3: 章节内精选")
    print("=" * 80)
    final_slides: dict[str, list[SlideInfo]] = {}

    for chapter in CHAPTERS:
        chapter_id = chapter["id"]
        slides = chapter_slides.get(chapter_id, [])

        if not slides:
            print(f"  章节 {chapter_id}: 无幻灯片")
            final_slides[chapter_id] = []
            continue

        # 计算每张幻灯片的得分
        scored_slides = []
        for slide_info in slides:
            score = _calculate_slide_score(slide_info, chapter["priority"])
            scored_slides.append((score, slide_info))

        # 按得分排序
        scored_slides.sort(key=lambda x: x[0], reverse=True)

        # 保留得分最高的（可根据需要调整数量）
        selected = [slide_info for _, slide_info in scored_slides]
        final_slides[chapter_id] = selected

        print(f"  章节 {chapter_id}: 保留 {len(selected)} 张")
        for score, slide_info in scored_slides[:3]:  # 显示前3张
            print(f"    - [{score}] {slide_info.slide_title[:50]}...")

    # 创建最终演示文稿
    print()
    print("=" * 80)
    print("步骤 4: 创建最终演示文稿")
    print("=" * 80)

    try:
        template_prs = Presentation(template_path)
    except Exception as e:
        logger.error(f"加载模板失败: {e}")
        return

    slide_count = 0
    integration_log: list[str] = []

    # 按章节顺序添加幻灯片
    for chapter in CHAPTERS:
        chapter_id = chapter["id"]
        chapter_title = chapter["title"]
        slides = final_slides.get(chapter_id, [])

        if not slides:
            continue

        # 添加章节过渡页
        add_chapter_divider_slide(template_prs, chapter_id, chapter_title)
        slide_count += 1
        integration_log.append(f"章节 {chapter_id} 过渡页已添加")

        # 添加该章节的幻灯片
        for slide_info in slides:
            try:
                new_slide = clone_slide_into_presentation(template_prs, slide_info.slide_obj, layout_index=0)
                standardize_slide_layout(new_slide)
                slide_count += 1

                log_msg = f"章节 {chapter_id}: 添加幻灯片 '{slide_info.slide_title[:50]}...' (来源: {slide_info.source_ppt}, 原始索引: {slide_info.source_index + 1})"  # noqa: E501
                integration_log.append(log_msg)
                print(f"  [OK] {log_msg}")
            except Exception as e:
                logger.error(f"添加幻灯片失败: {e}")
                integration_log.append(f"错误: 添加幻灯片失败 - {e}")

    # 保存最终文件
    output_pptx = os.path.join(output_dir, "Final_Report.pptx")
    try:
        template_prs.save(output_pptx)
        print(f"\n[OK] Final file saved: {output_pptx}")
    except Exception as e:
        logger.error(f"保存文件失败: {e}")
        return

    # 输出整合日志
    print()
    print("=" * 80)
    print("整合日志 (Integration Log)")
    print("=" * 80)
    for log_msg in integration_log:
        print(f"  {log_msg}")

    # 保存日志到文件
    log_file = os.path.join(output_dir, "integration_log.txt")
    try:
        with open(log_file, "w", encoding="utf-8") as f:
            f.write("=" * 80 + "\n")
            f.write("CSR 规范 PPT 整合日志\n")
            f.write("=" * 80 + "\n\n")
            f.write(f"输入文件数: {len(source_paths)}\n")
            f.write(f"原始幻灯片数: {len(slide_infos)}\n")
            f.write(f"去重后幻灯片数: {len(deduplicated_slides)}\n")
            f.write(f"最终幻灯片数: {slide_count}\n\n")
            f.write("=" * 80 + "\n")
            f.write("处理详情\n")
            f.write("=" * 80 + "\n\n")
            for log_msg in integration_log:
                f.write(f"{log_msg}\n")
        print(f"\n[OK] Integration log saved: {log_file}")
    except Exception as e:
        logger.error(f"保存日志文件失败: {e}")

    print()
    print("=" * 80)
    print("处理完成")
    print("=" * 80)
    print(f"[OK] Generated {slide_count} slides (including chapter dividers)")
    print(f"[OK] Final file: {output_pptx}")
    print(f"[OK] Integration log: {log_file}")


if __name__ == "__main__":
    integrate_ppts()
