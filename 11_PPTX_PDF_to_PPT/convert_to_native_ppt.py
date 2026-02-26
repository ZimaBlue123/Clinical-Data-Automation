# -*- coding: utf-8 -*-
"""
PPTX/PDF 转原生可编辑 PPT（全要素原位蒙版重建版）

优化特性：
1. 引入 tqdm 进度条，打破黑盒焦虑。
2. 全要素提取：不仅提取表格，同时提取正文、标题等普通文本。
3. 纯血原位替换 (PPTX)：直接修改原 PPT，保留所有已有可编辑元素，仅对图片进行 OCR 坐标映射和遮盖覆盖。
4. 绝对尺寸锁定 (PDF)：动态读取原始页面比例，拒绝图像拉伸压扁。
"""
from __future__ import annotations

import os
# 屏蔽底层报警
os.environ['PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK'] = 'True'
os.environ["GLOG_minloglevel"] = "2"

import argparse
import logging
import tempfile
from pathlib import Path

import pandas as pd
import fitz  # PyMuPDF
from PIL import Image
import io
from tqdm import tqdm  # 进度条神器

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)
# 优化终端日志，使其与进度条和睦相处
logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")


class ClinicalDocConverter:
    def __init__(self) -> None:
        try:
            from paddleocr.paddleocr import PPStructure
        except ImportError:
            try:
                from paddleocr import PPStructure
            except Exception as exc:
                raise ImportError("❌ 核心引擎导入失败！请检查 paddleocr 安装。") from exc

        logger.info("⚙️ 正在挂载 PP-Structure 全要素解析引擎...")
        self.table_engine = PPStructure(show_log=False, image_orientation=False, lang="ch")

    def process_file(self, input_file: Path, output_pptx: Path) -> None:
        ext = input_file.suffix.lower()

        with tempfile.TemporaryDirectory() as temp_dir:
            if ext == ".pdf":
                self._process_pdf_to_ppt(input_file, output_pptx, temp_dir)
            elif ext == ".pptx":
                self._process_pptx_inplace(input_file, output_pptx, temp_dir)
            else:
                raise ValueError("仅支持 .pdf 或 .pptx 格式。")

        logger.info(f"\n🎉 完美收工！成品已保存至: {output_pptx}")

    # ==========================================
    # 模式一：处理 PDF（新建等比例 PPT + 背景铺底 + 原位蒙版）
    # ==========================================
    def _process_pdf_to_ppt(self, pdf_path: Path, output_pptx: Path, temp_dir: str):
        logger.info("📄 检测到 PDF，正在进行页面等比光栅化...")
        doc = fitz.open(pdf_path)
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]

        # 动态锁定 PPT 尺寸（以第一页为准）
        page_0 = doc.load_page(0)
        prs.slide_width = Pt(page_0.rect.width)
        prs.slide_height = Pt(page_0.rect.height)

        for page_num in tqdm(range(len(doc)), desc="🖼️ PDF页面转换及解析"):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(matrix=fitz.Matrix(300/72, 300/72), alpha=False)
            img_path = os.path.join(temp_dir, f"page_{page_num}.png")
            pix.save(img_path)

            slide = prs.slides.add_slide(blank_layout)
            # 铺设底图
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            # 解析并覆盖文字/表格
            self._analyze_and_overlay(img_path, slide, 0, 0, prs.slide_width, prs.slide_height)

        prs.save(str(output_pptx))

    # ==========================================
    # 模式二：处理 PPTX（打开原文件 -> 找图片 -> 贴狗皮膏药遮盖）
    # ==========================================
    def _process_pptx_inplace(self, pptx_path: Path, output_pptx: Path, temp_dir: str):
        logger.info("📊 检测到 PPTX，启动原位无损替换模式...")
        prs = Presentation(str(pptx_path))
        
        # 收集所有图片对象，方便做进度条
        pic_tasks = []
        for s_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    pic_tasks.append((slide, shape, s_idx))

        if not pic_tasks:
            logger.warning("⚠️ 没在你的 PPT 里找到任何图片，直接保存原件。")
            prs.save(str(output_pptx))
            return

        for slide, shape, s_idx in tqdm(pic_tasks, desc="🔍 抓取图片并生成原生文本/表格"):
            img_blob = shape.image.blob
            img_ext = shape.image.ext
            img_path = os.path.join(temp_dir, f"temp_{s_idx}.{img_ext}")
            
            with open(img_path, "wb") as f:
                f.write(img_blob)
                
            # 获取图片在幻灯片上的绝对坐标和尺寸
            s_left, s_top = shape.left, shape.top
            s_width, s_height = shape.width, shape.height
            
            # 核心：将 AI 识别的数据直接盖在这张图片上方！
            self._analyze_and_overlay(img_path, slide, s_left, s_top, s_width, s_height)

        prs.save(str(output_pptx))

    # ==========================================
    # 核心引擎：解析图像 -> 计算相对坐标映射 -> 生成遮罩元素
    # ==========================================
    def _analyze_and_overlay(self, img_path: str, slide, base_left, base_top, base_width, base_height):
        # 1. 获取原图像素大小，用于计算缩放率
        with Image.open(img_path) as img:
            img_w_px, img_h_px = img.size
            
        scale_x = base_width / img_w_px
        scale_y = base_height / img_h_px

        # 2. 调用 AI 引擎识别版面
        result = self.table_engine(img_path)

        for region in result:
            bbox = region['bbox'] # [x1, y1, x2, y2]
            
            # 将图片的像素坐标，映射为 PPT 中的 Pt 坐标
            r_left = base_left + (bbox[0] * scale_x)
            r_top = base_top + (bbox[1] * scale_y)
            r_width = (bbox[2] - bbox[0]) * scale_x
            r_height = (bbox[3] - bbox[1]) * scale_y

            # ---- 处理表格 ----
            if region.get("type") == "table":
                try:
                    df_list = pd.read_html(region["res"]["html"])
                    if df_list:
                        df = df_list[0].fillna("").astype(str)
                        self._render_dataframe_to_slide(df, slide, r_left, r_top, r_width, r_height)
                except Exception as e:
                    pass

            # ---- 处理普通文字/标题 ----
            elif region.get("type") in ["text", "title", "figure_caption", "table_caption"]:
                try:
                    # 将识别出的多行文本拼合
                    texts = [item['text'] for item in region['res']]
                    full_text = "\n".join(texts)
                    self._render_text_to_slide(full_text, slide, r_left, r_top, r_width, r_height)
                except Exception as e:
                    pass

    # ==========================================
    # PPT 渲染组件：生成带白底的文本框（蒙版）
    # ==========================================
    def _render_text_to_slide(self, text: str, slide, left, top, width, height) -> None:
        txBox = slide.shapes.add_textbox(left, top, width, height)
        # 填充纯白色作为蒙版，遮住原图里的字
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Arial"
        p.font.size = Pt(10)

    # ==========================================
    # PPT 渲染组件：生成带白底的表格（蒙版）
    # ==========================================
    def _render_dataframe_to_slide(self, df: pd.DataFrame, slide, left, top, width, height) -> None:
        rows, cols = df.shape
        shape = slide.shapes.add_table(rows + 1, cols, left, top, width, height)
        table = shape.table

        # 填充纯白色表格底色
        for r_idx in range(rows + 1):
            for c_idx in range(cols):
                cell = table.cell(r_idx, c_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

        for col_idx, col_name in enumerate(df.columns):
            self._write_cell(table.cell(0, col_idx), str(col_name), is_header=True)

        for row_idx in range(rows):
            for col_idx in range(cols):
                val = df.iloc[row_idx, col_idx]
                self._write_cell(table.cell(row_idx + 1, col_idx), str(val))

    def _write_cell(self, cell, text: str, is_header: bool = False) -> None:
        cell.text = text
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Arial"
                if is_header:
                    run.font.size = Pt(10)
                    run.font.bold = True
                else:
                    run.font.size = Pt(9)

# ----------------- 启动区 -----------------
def main() -> None:
    base_dir = Path(__file__).resolve().parent
    input_dir = base_dir / "input"
    output_dir = base_dir / "output"
    
    input_dir.mkdir(parents=True, exist_ok=True)
    output_dir.mkdir(parents=True, exist_ok=True)

    candidates = list(input_dir.glob("*.pdf")) + list(input_dir.glob("*.pptx"))
    if not candidates:
        print("📭 input 目录下未找到文件！")
        return
        
    input_path = sorted(candidates)[0]
    output_path = output_dir / f"{input_path.stem}_editable.pptx"

    converter = ClinicalDocConverter()
    converter.process_file(input_path, output_path)

if __name__ == "__main__":
    main()