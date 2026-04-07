#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT 内嵌图直接去水印 (Vibe Ultimate Edition)

优化项：
1. 扩大覆盖面积，适配宽比例长条形 Logo。
2. 动态采样算法，严格在水印外侧安全区提取背景色，防止“采样污染”导致色差。
3. 采用中位数值提取颜色，免疫噪点。
"""
from __future__ import annotations

import argparse
import io
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageStat
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ================= 核心参数配置 =================
# 全页图判定：形状面积占幻灯片面积比例阈值
AREA_RATIO_THRESHOLD = 0.8

# 覆盖矩形（比例 + 像素上限）- 已大幅扩大火力覆盖范围
PATCH_WIDTH_RATIO = 0.1    # 覆盖最后 10% 宽度（解决长条文字 Logo）
PATCH_HEIGHT_RATIO = 0.04   # 覆盖最后 4% 高度
PATCH_MAX_WIDTH_PX = 400    # 最大宽度放宽至 400 像素
PATCH_MAX_HEIGHT_PX = 100   # 最大高度放宽至 100 像素
# ================================================

def get_script_dir() -> Path:
    return Path(__file__).resolve().parent

def slide_is_editable(slide) -> bool:
    """判定是否为可编辑页面：忽略图片及空占位符。"""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue
        # 忽略母版遗留的空占位符
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            if not shape.has_text_frame or not shape.text.strip():
                continue
        # 遇到其他实质性元素，跳过
        return True
    return False

def _sample_background_color(img: Image.Image, patch_w: int, patch_h: int) -> tuple:
    """
    动态安全区采样：基于最终计算的覆盖区大小，向左侧推移提取背景色
    """
    w, h = img.size
    
    # 采样点：覆盖区正左侧 30 像素，底部向上偏移一点点避免最底边的黑边
    # 这是一个 10x10 的采样块
    sample_x = max(0, w - patch_w - 30)
    sample_y = max(0, h - patch_h // 2)
    
    box = (sample_x - 5, max(0, sample_y - 5), sample_x + 5, min(h, sample_y + 5))
    region = img.crop(box).convert("RGBA")
    
    # 使用中位数 (Median) 而不是平均值 (Mean)，能有效忽略个别噪点，提取出最纯粹的背景底色
    stat = ImageStat.Stat(region)
    return tuple(int(c) for c in stat.median)

def patch_image_pillow(blob: bytes, content_type: str) -> bytes:
    """执行物理覆盖"""
    img = Image.open(io.BytesIO(blob)).convert("RGBA")
    w, h = img.size

    # 计算覆盖范围
    patch_w = min(max(1, int(w * PATCH_WIDTH_RATIO)), PATCH_MAX_WIDTH_PX)
    patch_h = min(max(1, int(h * PATCH_HEIGHT_RATIO)), PATCH_MAX_HEIGHT_PX)
    
    x0 = max(0, w - patch_w)
    y0 = max(0, h - patch_h)
    box = [x0, y0, w, h]

    # 先算尺寸，再向外延展安全区进行取色
    color = _sample_background_color(img, patch_w, patch_h)
    
    draw = ImageDraw.Draw(img)
    draw.rectangle(box, fill=color)

    # 封装备份
    buf = io.BytesIO()
    ct = content_type.lower()
    if "png" in ct:
        img.save(buf, format="PNG", compress_level=6)
    else:
        # JPEG 不支持 RGBA，转回 RGB
        img.convert("RGB").save(buf, format="JPEG", quality=100, subsampling=0)
    
    buf.seek(0)
    return buf.read()

def process_presentation(input_path: Path, output_path: Path) -> int:
    """遍历 PPT 执行重构"""
    prs = Presentation(str(input_path))
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    slide_area = float(slide_w * slide_h)
    patched_count = 0

    for slide in prs.slides:
        if slide_is_editable(slide):
            continue

        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            
            try:
                rId = shape._pic.blip_rId
            except Exception:
                continue
            if not rId:
                continue

            try:
                shape_area = float(shape.width * shape.height)
            except Exception:
                continue
            if slide_area <= 0 or shape_area / slide_area < AREA_RATIO_THRESHOLD:
                continue

            try:
                image_part = shape.part.related_part(rId)
                blob = image_part.blob
                content_type = image_part.content_type or "image/png"
            except Exception:
                continue

            if getattr(image_part, "ext", None) == "jpg":
                content_type = "image/jpeg"

            try:
                new_blob = patch_image_pillow(blob, content_type)
            except Exception as e:
                print(f"[!] 警告：图像处理异常 — {e}", file=sys.stderr)
                continue

            if new_blob != blob:
                image_part.blob = new_blob
                patched_count += 1

    prs.save(str(output_path))
    return patched_count

def main() -> None:
    parser = argparse.ArgumentParser(description="PPT 去水印 - Vibe Ultimate 版")
    parser.add_argument("input_pptx", type=str, nargs="?", default=None, help="输入 PPT 路径")
    parser.add_argument("-o", "--output", type=str, default=None, help="输出 PPT 路径")
    args = parser.parse_args()

    base = get_script_dir()
    input_dir = base / "input"
    output_dir = base / "output"
    output_dir.mkdir(parents=True, exist_ok=True)

    if args.input_pptx:
        input_path = Path(args.input_pptx)
        if not input_path.is_absolute():
            input_path = base / input_path
    else:
        pptx_list = sorted(input_dir.glob("*.pptx"))
        if not pptx_list:
            print("[X] 致命错误：未指定且未找到 .pptx 目标。", file=sys.stderr)
            sys.exit(1)
        input_path = pptx_list[0]
        print(f"[*] 锁定目标：{input_path.name}")

    if not input_path.exists():
        print(f"[X] 致命错误：路径不存在 {input_path}", file=sys.stderr)
        sys.exit(1)

    if args.output:
        output_path = Path(args.output)
        if not output_path.is_absolute():
            output_path = output_dir / output_path.name
    else:
        output_path = output_dir / f"{input_path.stem}_clean.pptx"

    print("[*] 正在执行深度覆盖...")
    n = process_presentation(input_path, output_path)
    print(f"[+] 任务完成！已输出：{output_path}（共拔除 {n} 个 Logo）")

if __name__ == "__main__":
    main()
