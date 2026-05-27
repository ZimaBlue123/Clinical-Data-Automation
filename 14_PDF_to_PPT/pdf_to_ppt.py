# -*- coding: utf-8 -*-
"""
将指定文件夹下的 PDF 批量转换为 PPTX（每页一页幻灯片，以图片形式嵌入）。

整合自 PDF to PPT-V4，默认路径为本项目下的「13_PDF_to_PPT/input」文件夹：
- 待转换 PDF 默认目录：13_PDF_to_PPT/input/
- 生成的 PPTX 默认输出到 13_PDF_to_PPT/output/ 文件夹

依赖: PyMuPDF (fitz), python-pptx, Pillow
安装: pip install pymupdf python-pptx pillow

用法:
  python pdf_to_ppt.py                    # 使用默认目录
  python pdf_to_ppt.py "D:\\其他\\PDF"   # 指定 PDF 目录（输出到同目录）
  python pdf_to_ppt.py "D:\\PDF" "D:\\PPT"  # 指定 PDF 目录与输出目录
"""

import sys
import traceback
from io import BytesIO
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

# ========== 配置 ==========
BASE = Path(__file__).resolve().parent
# 默认：待转换 PDF 与生成的 PPTX 使用模块下的 input/output 文件夹
PDF_DIR = BASE / "input"
OUTPUT_DIR = BASE / "output"

# 幻灯片尺寸（英寸）
SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
# 渲染目标 PPI
TARGET_PPI = 192
DPI_MIN = 72
DPI_MAX = 300
# ========================

SLIDE_PX_W = int(SLIDE_WIDTH_IN * TARGET_PPI)
SLIDE_PX_H = int(SLIDE_HEIGHT_IN * TARGET_PPI)


def render_page_to_png_bytes(page):
    """
    将 fitz.Page 渲染并返回 PNG 的 BytesIO（裁切为幻灯片纵横比并缩放到目标像素）。
    """
    rect = page.rect
    pdf_w_in = rect.width / 72.0
    pdf_h_in = rect.height / 72.0

    dpi_x = SLIDE_PX_W / pdf_w_in
    dpi_y = SLIDE_PX_H / pdf_h_in
    dpi = max(dpi_x, dpi_y)
    dpi = int(round(dpi))
    dpi = max(DPI_MIN, min(dpi, DPI_MAX))

    pix = page.get_pixmap(dpi=dpi, alpha=False)
    img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)

    img_ratio = img.width / img.height
    slide_ratio = SLIDE_PX_W / SLIDE_PX_H

    if img_ratio > slide_ratio:
        new_w = int(img.height * slide_ratio)
        left = (img.width - new_w) // 2
        img = img.crop((left, 0, left + new_w, img.height))
    else:
        new_h = int(img.width / slide_ratio)
        top = (img.height - new_h) // 2
        img = img.crop((0, top, img.width, top + new_h))

    img = img.resize((SLIDE_PX_W, SLIDE_PX_H), Image.LANCZOS)
    bio = BytesIO()
    img.save(bio, format="PNG", optimize=True)
    bio.seek(0)
    return bio


def pdf_to_ppt(pdf_path: Path, output_ppt_path: Path) -> None:
    """
    将单个 PDF 转为 pptx。若输出已存在且比 PDF 新则跳过。
    
    Args:
        pdf_path: PDF 文件路径
        output_ppt_path: 输出 PPTX 文件路径
        
    Raises:
        FileNotFoundError: PDF 文件不存在
        ValueError: PDF 文件无法打开或转换失败
    """
    if not pdf_path.exists():
        raise FileNotFoundError(f"PDF 文件不存在: {pdf_path}")
    
    doc = None
    try:
        # 检查输出文件是否已存在且为最新
        if output_ppt_path.exists():
            try:
                if output_ppt_path.stat().st_mtime >= pdf_path.stat().st_mtime:
                    print(f"跳过（已存在且为最新）: {output_ppt_path.name}")
                    return
            except OSError as e:
                print(f"警告: 无法比较文件时间戳: {e}")
        
        doc = fitz.open(pdf_path)
        if doc.page_count == 0:
            raise ValueError(f"PDF 文件没有页面: {pdf_path}")

        prs = Presentation()
        prs.slide_width = Inches(SLIDE_WIDTH_IN)
        prs.slide_height = Inches(SLIDE_HEIGHT_IN)

        total_pages = doc.page_count
        for idx in range(total_pages):
            try:
                page = doc.load_page(idx)
                img_stream = render_page_to_png_bytes(page)
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(
                    img_stream,
                    left=Inches(0),
                    top=Inches(0),
                    width=prs.slide_width,
                    height=prs.slide_height,
                )
                print(f"  {pdf_path.name}: 第 {idx + 1}/{total_pages} 页 完成")
            except Exception as e:
                print(f"警告: 处理第 {idx + 1} 页时出错: {e}")
                continue

        prs.save(output_ppt_path)
        print(f"✔ 已生成: {output_ppt_path}")
        
    except FileNotFoundError:
        raise
    except Exception as e:
        print(f"转换失败: {pdf_path}")
        traceback.print_exc()
        raise ValueError(f"PDF 转 PPT 失败: {e}") from e
    finally:
        if doc is not None:
            doc.close()


def main(target_dir: Path, out_dir: Path) -> None:
    """
    批量转换 PDF 文件为 PPTX。
    
    Args:
        target_dir: 包含 PDF 文件的目录
        out_dir: 输出目录
    """
    if not target_dir.exists():
        print(f"目标文件夹不存在: {target_dir}")
        return
    
    if not target_dir.is_dir():
        print(f"路径不是目录: {target_dir}")
        return

    try:
        out_dir.mkdir(parents=True, exist_ok=True)
    except OSError as e:
        print(f"无法创建输出目录 {out_dir}: {e}")
        return

    try:
        pdf_files = sorted([p for p in target_dir.iterdir() if p.is_file() and p.suffix.lower() == ".pdf"])
    except PermissionError as e:
        print(f"无法访问目录 {target_dir}: {e}")
        return

    if not pdf_files:
        print(f"未在 {target_dir} 找到 PDF 文件。")
        return

    print(f"找到 {len(pdf_files)} 个 PDF，输出目录：{out_dir}\n")
    success_count = 0
    for pdf_path in pdf_files:
        output_ppt_path = out_dir / (pdf_path.stem + ".pptx")
        try:
            pdf_to_ppt(pdf_path, output_ppt_path)
            success_count += 1
        except FileNotFoundError as e:
            print(f"文件未找到: {e}")
        except ValueError as e:
            print(f"转换错误: {e}")
        except Exception as e:
            print(f"处理文件时出现未捕捉异常: {pdf_path}")
            traceback.print_exc()
    
    print(f"\n全部处理完成。成功转换 {success_count}/{len(pdf_files)} 个文件。")


if __name__ == "__main__":
    pdf_dir = PDF_DIR
    out_dir = OUTPUT_DIR
    if len(sys.argv) >= 2:
        pdf_dir = Path(sys.argv[1])
        out_dir = Path(sys.argv[2]) if len(sys.argv) >= 3 else pdf_dir
    main(pdf_dir, out_dir)
